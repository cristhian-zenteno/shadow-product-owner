"""
Generate sample documents (PDF, DOCX, PPTX) for testing Component D.

These fixtures are used to test document-to-Markdown conversion via MarkItDown.
"""

from pathlib import Path


def generate_pdf():
    """Generate a sample PDF using reportlab."""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
    except ImportError:
        print("⚠️  reportlab not installed. Run: uv add reportlab")
        return False
    
    pdf_path = Path(__file__).parent / "sample.pdf"
    
    c = canvas.Canvas(str(pdf_path), pagesize=letter)
    c.drawString(100, 750, "Sample Product Requirements Document")
    c.drawString(100, 720, "")
    c.drawString(100, 690, "Feature: One-Click Checkout")
    c.drawString(100, 660, "")
    c.drawString(100, 630, "User Story:")
    c.drawString(100, 600, "As a customer, I want to complete my purchase with a single click")
    c.drawString(100, 570, "so that I can checkout faster and reduce cart abandonment.")
    c.save()
    
    print(f"✓ Created: {pdf_path}")
    return True


def generate_docx():
    """Generate a sample DOCX using python-docx."""
    try:
        from docx import Document
    except ImportError:
        print("⚠️  python-docx not installed. Run: uv add python-docx")
        return False
    
    docx_path = Path(__file__).parent / "sample.docx"
    
    doc = Document()
    doc.add_heading("Sample Product Requirements Document", level=1)
    doc.add_heading("Feature: One-Click Checkout", level=2)
    doc.add_paragraph("User Story:")
    doc.add_paragraph(
        "As a customer, I want to complete my purchase with a single click "
        "so that I can checkout faster and reduce cart abandonment."
    )
    doc.add_heading("Acceptance Criteria", level=2)
    doc.add_paragraph("• The checkout button is visible on the cart page")
    doc.add_paragraph("• Clicking the button completes the purchase immediately")
    doc.add_paragraph("• User receives confirmation email within 1 minute")
    
    doc.save(str(docx_path))
    
    print(f"✓ Created: {docx_path}")
    return True


def generate_pptx():
    """Generate a sample PPTX using python-pptx."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("⚠️  python-pptx not installed. Run: uv add python-pptx")
        return False
    
    pptx_path = Path(__file__).parent / "sample.pptx"
    
    prs = Presentation()
    
    # Slide 1: Title
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = "One-Click Checkout"
    subtitle.text = "Product Requirements Presentation"
    
    # Slide 2: User Story
    bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = bullet_slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "User Story"
    tf = body_shape.text_frame
    tf.text = "As a customer, I want to complete my purchase with a single click"
    
    p = tf.add_paragraph()
    p.text = "So that I can checkout faster and reduce cart abandonment"
    
    prs.save(str(pptx_path))
    
    print(f"✓ Created: {pptx_path}")
    return True


if __name__ == "__main__":
    print("Generating sample documents for Component D tests...\n")
    
    results = []
    results.append(("PDF", generate_pdf()))
    results.append(("DOCX", generate_docx()))
    results.append(("PPTX", generate_pptx()))
    
    print("\n" + "=" * 60)
    success_count = sum(1 for _, success in results if success)
    print(f"Generated {success_count}/3 document fixtures")
    
    if success_count < 3:
        print("\nTo install missing dependencies:")
        if not results[0][1]:
            print("  uv add reportlab")
        if not results[1][1]:
            print("  uv add python-docx")
        if not results[2][1]:
            print("  uv add python-pptx")
